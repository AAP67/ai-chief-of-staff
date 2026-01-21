"""
File processing utilities for AI Chief of Staff
Handles PDF, PPTX, XLSX, CSV, TXT, and MD files
"""

import io
from typing import Union
import streamlit as st

# Import document processing libraries
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
    import pandas as pd
except ImportError:
    openpyxl = None
    pd = None


def process_uploaded_file(file) -> str:
    """
    Process an uploaded file and extract text content
    
    Args:
        file: Streamlit UploadedFile object
    
    Returns:
        str: Extracted text content from the file
    """
    file_extension = file.name.split('.')[-1].lower()
    
    try:
        if file_extension == 'pdf':
            return process_pdf(file)
        elif file_extension == 'pptx':
            return process_pptx(file)
        elif file_extension in ['xlsx', 'xls']:
            return process_excel(file)
        elif file_extension == 'csv':
            return process_csv(file)
        elif file_extension in ['txt', 'md']:
            return process_text(file)
        else:
            return f"[Unsupported file type: {file_extension}]"
    except Exception as e:
        return f"[Error processing {file.name}: {str(e)}]"


def process_pdf(file) -> str:
    """Extract text from PDF file"""
    if PdfReader is None:
        return "[PDF processing not available - PyPDF2 not installed]"
    
    try:
        pdf_reader = PdfReader(io.BytesIO(file.read()))
        text = ""
        
        for page_num, page in enumerate(pdf_reader.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text += f"\n--- Page {page_num} ---\n{page_text}\n"
        
        return text.strip() if text.strip() else "[PDF contains no extractable text]"
    except Exception as e:
        return f"[Error reading PDF: {str(e)}]"


def process_pptx(file) -> str:
    """Extract text from PowerPoint file"""
    if Presentation is None:
        return "[PowerPoint processing not available - python-pptx not installed]"
    
    try:
        prs = Presentation(io.BytesIO(file.read()))
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += f"{shape.text}\n"
        
        return text.strip() if text.strip() else "[PowerPoint contains no extractable text]"
    except Exception as e:
        return f"[Error reading PowerPoint: {str(e)}]"


def process_excel(file) -> str:
    """Extract data from Excel file"""
    if pd is None or openpyxl is None:
        return "[Excel processing not available - pandas/openpyxl not installed]"
    
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(io.BytesIO(file.read()))
        text = ""
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            text += f"\n--- Sheet: {sheet_name} ---\n"
            text += df.to_string(index=False)
            text += "\n\n"
        
        return text.strip() if text.strip() else "[Excel file contains no data]"
    except Exception as e:
        return f"[Error reading Excel: {str(e)}]"


def process_csv(file) -> str:
    """Extract data from CSV file"""
    if pd is None:
        return "[CSV processing not available - pandas not installed]"
    
    try:
        df = pd.read_csv(io.BytesIO(file.read()))
        text = f"CSV File: {file.name}\n\n"
        text += df.to_string(index=False)
        return text
    except Exception as e:
        return f"[Error reading CSV: {str(e)}]"


def process_text(file) -> str:
    """Extract text from text/markdown file"""
    try:
        content = file.read()
        # Try to decode as UTF-8, fall back to latin-1 if that fails
        try:
            text = content.decode('utf-8')
        except UnicodeDecodeError:
            text = content.decode('latin-1')
        
        return text.strip() if text.strip() else "[File is empty]"
    except Exception as e:
        return f"[Error reading text file: {str(e)}]"


def get_file_summary(file_data: dict) -> str:
    """
    Generate a brief summary of processed file content
    
    Args:
        file_data: Dictionary with 'name' and 'content' keys
    
    Returns:
        str: Summary of file content
    """
    content = file_data['content']
    name = file_data['name']
    
    # Count words and lines
    words = len(content.split())
    lines = len(content.split('\n'))
    
    return f"📄 {name} ({words} words, {lines} lines)"
